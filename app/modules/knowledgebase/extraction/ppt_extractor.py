"""
Structured PPT/PPTX extraction with slide titles, text, tables, notes, and images.
"""

from __future__ import annotations

import hashlib
import logging
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.modules.knowledgebase.extraction.models import StructuredDocument
from app.modules.knowledgebase.extraction.table_utils import table_rows_to_markdown
from app.modules.knowledgebase.extraction.vision import describe_image_bytes, reset_vision_session

logger = logging.getLogger(__name__)


def _extract_shape_table(shape) -> str | None:
    table = shape.table
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return table_rows_to_markdown(rows)


def _extract_slide_images(slide, seen_hashes: set[str]) -> list[str]:
    descriptions: list[str] = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            image = shape.image
            image_bytes = image.blob
            if not image_bytes:
                continue
            image_hash = hashlib.md5(image_bytes).hexdigest()
            if image_hash in seen_hashes:
                continue
            seen_hashes.add(image_hash)

            content_type = image.content_type or "image/png"
            description = describe_image_bytes(image_bytes, content_type)
            if description:
                descriptions.append(description)
        except Exception:
            logger.exception("Failed to extract/describe PPT image")
            continue
    return descriptions


def extract_pptx_structured(file_path: Path) -> StructuredDocument:
    """Extract structured content from a PPT or PPTX presentation."""
    document = StructuredDocument()
    reset_vision_session()
    presentation = Presentation(file_path)
    seen_image_hashes: set[str] = set()

    for slide_index, slide in enumerate(presentation.slides, start=1):
        document.add("slide_marker", f"--- Slide {slide_index} ---")

        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                document.add("slide_title", slide.shapes.title.text.strip())

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_markdown = _extract_shape_table(shape)
                    if table_markdown:
                        document.add(
                            "table",
                            table_markdown,
                            metadata={"slide": slide_index},
                        )
                    continue

                if not getattr(shape, "has_text_frame", False):
                    continue
                if shape == slide.shapes.title:
                    continue

                text = shape.text.strip()
                if text:
                    document.add(
                        "paragraph",
                        text,
                        metadata={"slide": slide_index},
                    )

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    document.add(
                        "speaker_notes",
                        notes,
                        metadata={"slide": slide_index},
                    )

            for description in _extract_slide_images(slide, seen_image_hashes):
                document.add(
                    "image_description",
                    description,
                    metadata={"slide": slide_index},
                )
        except Exception:
            logger.exception(
                "PPT slide extraction failed file=%s slide=%s",
                file_path,
                slide_index,
            )
            continue

    return document
